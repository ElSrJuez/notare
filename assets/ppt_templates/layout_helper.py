"""Utility to list slide layout names in the first .pptx file found in the template folder.
Run:  python layout_helper.py
"""
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE as PPT
from pathlib import Path
import sys

# This helper now only lists the layouts and their placeholder info.
# All interactive mapping and JSON persistence has been removed
# because the backend no longer uses `layout_map.json`.

TEMPLATE_DIR = Path(__file__).resolve().parent
DEFAULTS = [
    "Title Slide",
    "Title and Content",
    "Section Header",
    "Two Content",
    "Comparison",
    "Title Only",
    "Blank",
    "Content with Caption",
    "Picture with Caption",
]

def main() -> None:
    try:
        template_path = next(f for f in TEMPLATE_DIR.glob("*.pptx"))
    except StopIteration:
        print("No .pptx template found in", TEMPLATE_DIR)
        sys.exit(1)

    prs = Presentation(str(template_path))
    print(f"Layouts in {template_path.name} (index | name | title? | body?)")
    for idx, layout in enumerate(prs.slide_layouts):
        has_title = any(ph.placeholder_format.type in (PPT.TITLE, PPT.CENTER_TITLE, PPT.SUBTITLE) for ph in layout.placeholders)
        has_body = any(ph.placeholder_format.type == PPT.BODY for ph in layout.placeholders)
        plh_types = ", ".join({PPT(ph.placeholder_format.type).name for ph in layout.placeholders})
        print(f" {idx:2}: {layout.name} | title={has_title} | body={has_body} | placeholders: {plh_types}")

    # Done.

if __name__ == "__main__":
    main()
