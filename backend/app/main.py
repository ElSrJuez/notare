from fastapi import FastAPI, HTTPException, UploadFile, File, Form, Depends
from pydantic import BaseModel, HttpUrl
import httpx
from bs4 import BeautifulSoup
from fastapi.middleware.cors import CORSMiddleware
from readability import Document
import logging
from typing import List, Optional, Annotated
from markdownify import MarkdownConverter
import os, io, json
import openai
from pptx import Presentation
from pptx.util import Inches, Pt
from fastapi.responses import StreamingResponse, FileResponse
from .llm_provider import OpenAIProvider, LlamaHTTPProvider, BaseProvider, OutlineError
from pathlib import Path
from .config import load_config
import tempfile  # added
from fastapi.staticfiles import StaticFiles
from pathlib import Path as _Path

# ---------------------------------------------------------------------------
# Default layout names (internal use)
# ---------------------------------------------------------------------------

_TITLE_LAYOUT_NAME: str = "Title Slide"
_CONTENT_LAYOUT_NAME: str = "Title and Content"

app = FastAPI(title="Notāre Backend")
cfg = load_config()

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("notare.pptx")
logger.setLevel(logging.INFO)

# Add CORS middleware to allow frontend requests.
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=["X-Notare-Template-Diagnostics"],
)

class NormalizeRequest(BaseModel):
    url: HttpUrl

@app.post("/normalize")
async def normalize_page(req: NormalizeRequest):
    """Fetches the web page, strips boilerplate, and returns simplified HTML."""
    try:
        async with httpx.AsyncClient() as client:
            resp = await client.get(str(req.url), timeout=10)
            resp.raise_for_status()
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))

    doc = Document(resp.text, url=str(req.url))
    summary_html = doc.summary(html_partial=True)

    # Further strip unwanted tags but keep structure
    soup = BeautifulSoup(summary_html, "html.parser")
    for tag in soup(["script", "style", "form"]):
        tag.decompose()

    # Debug: log remaining blockquotes or <br> that could break highlighting
    for bq in soup.find_all("blockquote"):
        logging.info("BLOCKQUOTE remaining: %s", bq.get_text(strip=True)[:120])
    if soup.find("br"):
        logging.info("<br> tags remain in document – may split sentences")

    clean_html = str(soup)

    return {"clean_html": clean_html}

class OutlineRequest(BaseModel):
    """Annotated HTML sent from the frontend to generate an outline."""
    html: str

# Custom converter to mark <mark class="notare-mark">
class MarkAwareConverter(MarkdownConverter):
    def convert_mark(self, el, text, *args, **kwargs):
        if 'notare-mark' in el.get('class', []):
            return f"<<mark>>{text}<</mark>>"
        return text

def html_to_markdown(html: str) -> str:
    conv = MarkAwareConverter(bullets='*')
    return conv.convert(html)

@app.post("/outline")
async def outline(req: OutlineRequest):
    """Return markdown with custom highlight markers for LLM processing."""
    markdown = html_to_markdown(req.html)
    return {"markdown": markdown}

class PPTXRequest(BaseModel):
    html: str

def validate_template(prs):
    """Return structured diagnostics for required layouts and placeholders.

    Output Example::
        [
            {"layout": "Title Slide", "status": "error", "message": "Layout not found"},
            {"layout": "Title and Content", "status": "warning", "message": "Missing BODY placeholder"}
        ]
    """

    diagnostics: list[dict] = []

    def append(layout: str, status: str, message: str):
        diagnostics.append({"layout": layout, "status": status, "message": message})

    def get_plh_types(layout):
        return {ph.placeholder_format.type for ph in layout.placeholders}

    # --- Title Slide ---
    title_slide = next((l for l in prs.slide_layouts if l.name == _TITLE_LAYOUT_NAME), None)
    if not title_slide:
        append("Title Slide", "error", "Layout not found")
    else:
        types = get_plh_types(title_slide)
        if 1 not in types:  # 1 == TITLE
            append("Title Slide", "warning", "TITLE placeholder missing")
        else:
            append("Title Slide", "ok", "Layout present with TITLE placeholder")

    # --- Title and Content ---
    title_content = next((l for l in prs.slide_layouts if l.name == _CONTENT_LAYOUT_NAME), None)
    if not title_content:
        append("Title and Content", "error", "Layout not found")
    else:
        types = get_plh_types(title_content)
        missing: list[str] = []
        if 1 not in types:
            missing.append("TITLE")
        if 2 not in types:  # 2 == BODY
            missing.append("BODY")
        if missing:
            append("Title and Content", "warning", f"Missing placeholder(s): {', '.join(missing)}")
        else:
            append("Title and Content", "ok", "Layout present with TITLE and BODY placeholders")

    return diagnostics

class SettingsPayload(BaseModel):
    provider: str = "openai"
    api_key: str | None = None
    model: Optional[str] = None
    endpoint: Optional[str] = None
    api_version: Optional[str] = None

    def build_provider(self) -> BaseProvider:
        name = self.provider.lower()
        if name in ("openai", "azure"):
            if not self.api_key and name != "llama":
                raise HTTPException(status_code=400, detail="api_key required for OpenAI/Azure provider")
            if name == "azure" and not self.endpoint:
                raise HTTPException(status_code=400, detail="endpoint required for Azure provider")
            return OpenAIProvider({"api_key": self.api_key, "model": self.model or "gpt-4o-chat-bison", "endpoint": self.endpoint, "api_version": self.api_version})
        if name == "llama":
            return LlamaHTTPProvider(base_url=self.endpoint or "http://localhost:8080/completions", model=self.model)
        raise HTTPException(status_code=400, detail=f"Unsupported provider '{self.provider}'")


@app.get("/healthz")
async def healthz():
    return {"ok": True}

# Fallback status endpoint (distinct path to avoid masking GUI)
@app.get("/__status")
async def status():
    return {"ok": True}

@app.post("/pptx")
async def generate_pptx(
    settings: Annotated[str, Form(...)],
    html: Annotated[str, Form(..., description="Annotated HTML from client")],
    template: Annotated[UploadFile | None, File(description="Optional PPTX template", alias="template")] = None,
    # layout_map upload field removed (legacy)
):
    """Generate a PPTX from annotated HTML via LLM outline.

    Expects multipart/form-data with:
        settings: JSON string with provider/api_key/model (required)
        template: .pptx file ≤5 MiB (optional)
        layout_map: .json mapping file (optional)
        html: annotated article HTML (required)
    """

    try:
        settings_obj = SettingsPayload.model_validate_json(settings)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Invalid settings JSON: {e}")

    if template and template.size and template.size > 5 * 1024 * 1024:
        raise HTTPException(status_code=400, detail="Template too large (limit 5 MiB)")

    provider = settings_obj.build_provider()

    markdown = html_to_markdown(html)

    try:
        outline = await provider.generate_outline(markdown)
    except OutlineError as e:
        raise HTTPException(status_code=500, detail=str(e))

    # Load template: uploaded > configured dir > default blank.
    prs_template = None
    cleanup_files = []
    if template is not None:
        tmp_path = Path(tempfile.gettempdir()) / f"notare_{template.filename}"
        with tmp_path.open("wb") as f:
            f.write(await template.read())
        prs_template = Presentation(str(tmp_path))
        cleanup_files.append(tmp_path)
    else:
        prs_template = Presentation()

    prs = prs_template
    diagnostics = validate_template(prs)
    logger.info("Template diagnostics: %s", diagnostics)

    def get_layout(name:str, fallback_idx:int):
        return next((l for l in prs.slide_layouts if l.name == name), prs.slide_layouts[fallback_idx])

    title_layout   = get_layout(_TITLE_LAYOUT_NAME, 0)
    content_layout = get_layout(_CONTENT_LAYOUT_NAME, 1)

    for idx, slide_data in enumerate(outline.get("slides", [])):
        layout = title_layout if idx == 0 else content_layout
        s = prs.slides.add_slide(layout)

        # Title
        title_shape = s.shapes.title
        if title_shape:
            title_shape.text = slide_data.get("title", "")
        else:  # fallback create title textbox
            s.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1)).text_frame.text = slide_data.get("title", "")

        # Bullets
        if idx != 0:
            body_ph = next((ph for ph in s.placeholders if ph.placeholder_format.type == 2), None)
            if body_ph:
                tf = body_ph.text_frame
            else:
                tf = s.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5)).text_frame
            tf.clear()
            for bullet in slide_data.get("bullets", []):
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = Pt(18)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    response = StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=outline.pptx"}
    )
    # Legacy flat header for backward compatibility (semicolon delimited key=value)
    def header_token(d: dict[str, str]):
        return f"{d['layout']}={d['status']}"

    header_val = "|".join(header_token(d) for d in diagnostics)
    response.headers["X-Notare-Template-Diagnostics"] = header_val
    # After streaming ensure temp files removed
    for p in cleanup_files:
        try:
            p.unlink(missing_ok=True)
        except Exception:
            pass
    return response

# ---------------------------------------------------------------------------
# Template-only validation endpoint (no PPTX generation)
# ---------------------------------------------------------------------------

@app.post("/template/validate")
async def validate_only(
    template: Annotated[UploadFile | None, File(description="PPTX template", alias="template")] = None,
):
    """Validate a PPTX template and return structured diagnostics.

    This endpoint allows the frontend to upload a template, get precise
    validation feedback, and decide whether to continue with slide generation.
    """

    # Load presentation – from uploaded file or configured template directory
    prs: Presentation | None = None
    cleanup_path: Path | None = None

    if template is not None:
        if template.size and template.size > 5 * 1024 * 1024:
            raise HTTPException(status_code=400, detail="Template too large (limit 5 MiB)")

        tmp_path = Path(tempfile.gettempdir()) / f"notare_validate_{template.filename}"
        with tmp_path.open("wb") as f:
            f.write(await template.read())
        prs = Presentation(str(tmp_path))
        cleanup_path = tmp_path
    else:
        prs = Presentation()

    diagnostics = validate_template(prs)

    # Determine overall state
    has_error = any(d["status"] == "error" for d in diagnostics)
    has_warning = any(d["status"] == "warning" for d in diagnostics)

    summary = "error" if has_error else "warning" if has_warning else "ok"

    # Clean up temp file
    if cleanup_path:
        try:
            cleanup_path.unlink(missing_ok=True)
        except Exception:
            pass

    return {"summary": summary, "diagnostics": diagnostics}

# Mount frontend build (if present) as static site
_frontend_path = _Path(__file__).resolve().parents[2] / "frontend" / "dist"
if _frontend_path.exists():
    app.mount("/", StaticFiles(directory=str(_frontend_path), html=True), name="frontend")
